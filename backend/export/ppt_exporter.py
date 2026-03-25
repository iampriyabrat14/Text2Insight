"""
PowerPoint (.pptx) exporter.

Slide layout per Q&A pair:
  Slide A : Question + AI Summary  (existing design)
  Slide B : Data Table + Bar Chart (new — only when result data is present)

Bookend slides:
  Slide 1      : Cover
  Last slide   : End card
"""
import io
from datetime import datetime, timezone

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from backend.export.models import ExportSession
from backend.export.pdf_exporter import _pick_chart_cols

_BLUE   = RGBColor(0x25, 0x63, 0xEB)
_INDIGO = RGBColor(0x4F, 0x46, 0xE5)
_DARK   = RGBColor(0x0F, 0x17, 0x2A)
_SLATE  = RGBColor(0x33, 0x41, 0x55)
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT  = RGBColor(0xF8, 0xFA, 0xFF)
_MUTED  = RGBColor(0x94, 0xA3, 0xB8)
_ACCENT = RGBColor(0xEF, 0xF6, 0xFF)
_TH_BG  = RGBColor(0x1E, 0x3A, 0x8A)
_ALT_BG = RGBColor(0xF1, 0xF5, 0xF9)

_W = Inches(13.33)
_H = Inches(7.5)


# ── Low-level shape helpers ─────────────────────────────────────────────────

def _new_prs():
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def _text(slide, l, t, w, h, txt, size=14, bold=False, color=_DARK,
          align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return tb


# ── Slide builders ─────────────────────────────────────────────────────────

def _slide_cover(prs, session):
    slide = _blank(prs)
    _rect(slide, 0, 0, _W, _H,   fill=RGBColor(0x08, 0x0E, 0x1D))
    _rect(slide, 0, 0, _W, Inches(4), fill=RGBColor(0x0F, 0x1A, 0x3A))
    _rect(slide, 0, 0, Inches(0.18), _H, fill=_BLUE)

    _text(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(1.2),
          "AI Insights Report", size=40, bold=True, color=_WHITE, font="Calibri")
    _text(slide, Inches(0.6), Inches(2.9), Inches(10), Inches(0.7),
          session.title, size=18, color=RGBColor(0xBB, 0xDE, 0xFB))
    _rect(slide, Inches(0.6), Inches(3.75), Inches(4), Inches(0.04), fill=_BLUE)

    export_dt = datetime.now(timezone.utc).strftime("%B %d, %Y")
    _text(slide, Inches(0.6), Inches(3.95), Inches(10), Inches(0.5),
          f"{session.username}  ·  {export_dt}  ·  {len(session.qa_pairs)} questions",
          size=12, color=_MUTED)


def _slide_qa(prs, idx, pair):
    slide = _blank(prs)
    _rect(slide, 0, 0, _W, _H,         fill=RGBColor(0xF8, 0xFA, 0xFF))
    _rect(slide, 0, 0, _W, Inches(0.08), fill=_BLUE)

    # Question badge
    _rect(slide, Inches(0.4), Inches(0.25), Inches(0.85), Inches(0.5), fill=_BLUE)
    _text(slide, Inches(0.4), Inches(0.25), Inches(0.85), Inches(0.5),
          f"Q {idx}", size=14, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    # Left: question box
    _rect(slide, Inches(0.4), Inches(0.9), Inches(5.6), Inches(5.8),
          fill=_ACCENT, line=RGBColor(0xBF, 0xDB, 0xFE), line_w=Pt(1))
    _text(slide, Inches(0.52), Inches(0.95), Inches(5.3), Inches(0.4),
          "QUESTION", size=8, bold=True, color=_BLUE, font="Calibri")
    _text(slide, Inches(0.52), Inches(1.35), Inches(5.3), Inches(5.2),
          pair["query"], size=15, bold=True, color=_DARK)

    # Right: summary box
    _rect(slide, Inches(6.4), Inches(0.9), Inches(6.5), Inches(5.8),
          fill=_WHITE, line=RGBColor(0xE2, 0xE8, 0xF0), line_w=Pt(1))
    _text(slide, Inches(6.55), Inches(0.95), Inches(6.2), Inches(0.4),
          "AI SUMMARY", size=8, bold=True, color=_INDIGO, font="Calibri")
    _text(slide, Inches(6.55), Inches(1.4), Inches(6.2), Inches(5.1),
          pair["summary"], size=13, color=_SLATE)

    _text(slide, _W - Inches(1), _H - Inches(0.4), Inches(0.8), Inches(0.35),
          str(idx + 1), size=9, color=_MUTED, align=PP_ALIGN.RIGHT)


def _slide_data(prs, idx, pair):
    """Data slide: table on the left, bar chart on the right."""
    result  = pair.get("result_data") or {}
    columns = result.get("columns", [])
    rows    = result.get("rows",    [])
    if not columns or not rows:
        return

    slide = _blank(prs)
    _rect(slide, 0, 0, _W, _H,           fill=RGBColor(0xF8, 0xFA, 0xFF))
    _rect(slide, 0, 0, _W, Inches(0.08), fill=_INDIGO)

    # Header
    _rect(slide, Inches(0.4), Inches(0.18), Inches(0.85), Inches(0.5), fill=_INDIGO)
    _text(slide, Inches(0.4), Inches(0.18), Inches(0.85), Inches(0.5),
          f"Q {idx}", size=14, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.4), Inches(0.22), Inches(8), Inches(0.4),
          "Data Table & Chart", size=13, bold=True, color=_DARK)

    # ── Table (left half) ─────────────────────────────────────────────────
    max_tbl_rows = min(len(rows), 14)
    max_tbl_cols = min(len(columns), 6)
    disp_cols    = columns[:max_tbl_cols]
    disp_rows    = rows[:max_tbl_rows]

    tbl_shape = slide.shapes.add_table(
        max_tbl_rows + 1, max_tbl_cols,
        Inches(0.3), Inches(0.85),
        Inches(5.9), Inches(6.4),
    )
    tbl = tbl_shape.table

    # Column widths
    col_w = Inches(5.9) // max_tbl_cols
    for j in range(max_tbl_cols):
        tbl.columns[j].width = col_w

    # Header row
    for j, col in enumerate(disp_cols):
        cell = tbl.cell(0, j)
        cell.text = col[:18]
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold      = True
        run.font.size      = Pt(8)
        run.font.color.rgb = _WHITE
        fill = cell._tc.get_or_add_tcPr()
        from pptx.oxml.ns import qn as _qn
        from lxml import etree
        solidFill = etree.SubElement(fill, _qn("a:solidFill"))
        srgb      = etree.SubElement(solidFill, _qn("a:srgbClr"))
        srgb.set("val", "1E3A8A")

    # Data rows
    for i, row in enumerate(disp_rows, start=1):
        is_alt = (i % 2 == 0)
        for j, col in enumerate(disp_cols):
            val  = row.get(col)
            s    = str(val) if val is not None else "—"
            cell = tbl.cell(i, j)
            cell.text = s[:22]
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size      = Pt(7.5)
            run.font.color.rgb = _SLATE
            if is_alt:
                fill = cell._tc.get_or_add_tcPr()
                solidFill = etree.SubElement(fill, _qn("a:solidFill"))
                srgb      = etree.SubElement(solidFill, _qn("a:srgbClr"))
                srgb.set("val", "F1F5F9")

    # ── Chart (right half) ───────────────────────────────────────────────
    label_col, value_col = _pick_chart_cols(columns, rows)
    if label_col and value_col:
        chart_rows = rows[:12]
        try:
            labels = [str(r.get(label_col, ""))[:16] for r in chart_rows]
            values = tuple(
                float(str(r.get(value_col, 0)).replace(",", "").replace("$", "") or 0)
                for r in chart_rows
            )

            cd = ChartData()
            cd.categories = labels
            cd.add_series(value_col[:20], values)

            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(6.4), Inches(0.85),
                Inches(6.6), Inches(6.4),
                cd,
            )
            c = chart_shape.chart
            c.has_legend = False
            c.has_title  = True
            c.chart_title.text_frame.text = value_col[:30]
            c.chart_title.text_frame.paragraphs[0].runs[0].font.size  = Pt(11)
            c.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _SLATE

            # Style bars
            series = c.series[0]
            fill   = series.format.fill
            fill.solid()
            fill.fore_color.rgb = _BLUE

            # Style plot area
            c.plot_area.format.fill.background()
            c.chart_area.format.fill.background()

        except Exception:
            pass   # skip chart if data conversion fails


def _slide_end(prs, session):
    slide = _blank(prs)
    _rect(slide, 0, 0, _W, _H,           fill=RGBColor(0x08, 0x0E, 0x1D))
    _rect(slide, 0, 0, Inches(0.18), _H, fill=_INDIGO)
    _text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
          "End of Report", size=36, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
          f"{len(session.qa_pairs)} AI-generated summaries  ·  SQL-to-NLP",
          size=14, color=_MUTED, align=PP_ALIGN.CENTER)


# ── Public API ─────────────────────────────────────────────────────────────

def generate_ppt(session: ExportSession) -> bytes:
    prs = _new_prs()
    _slide_cover(prs, session)
    for idx, pair in enumerate(session.qa_pairs, start=1):
        _slide_qa(prs, idx, pair)
        _slide_data(prs, idx, pair)    # data slide (skipped if no rows)
    _slide_end(prs, session)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
